# src/office_ops/ppt_processor/content_models.py
import re
from pptx.enum.shapes import MSO_SHAPE_TYPE, MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from typing import List, Dict, Optional, Tuple
from config.loader import ConfigLoader  # 导入配置加载器

# ------------------------------ 静态工具方法 ------------------------------
def calculate_iou(box1: Tuple[float, float, float, float],
                  box2: Tuple[float, float, float, float]) -> float:
    """
    计算两个矩形的交并比（IOU）
    box格式：(left, top, width, height) 单位cm
    """
    # 转换为 (x1, y1, x2, y2) 坐标
    b1_x1, b1_y1, b1_w, b1_h = box1
    b1_x2 = b1_x1 + b1_w
    b1_y2 = b1_y1 + b1_h

    b2_x1, b2_y1, b2_w, b2_h = box2
    b2_x2 = b2_x1 + b2_w
    b2_y2 = b2_y1 + b2_h

    # 计算交集面积
    inter_x1 = max(b1_x1, b2_x1)
    inter_y1 = max(b1_y1, b2_y1)
    inter_x2 = min(b1_x2, b2_x2)
    inter_y2 = min(b1_y2, b2_y2)
    inter_area = max(0, inter_x2 - inter_x1) * max(0, inter_y2 - inter_y1)

    # 计算并集面积
    b1_area = b1_w * b1_h
    b2_area = b2_w * b2_h
    union_area = b1_area + b2_area - inter_area

    return inter_area / union_area if union_area != 0 else 0.0


def is_empty(s) -> bool:
    """严格判断字符串是否为空"""
    if s is None:
        return True
    if not isinstance(s, str):
        return False
    return s.strip() == ""

class UnitConverter:
    """单位转换工具类（EMU 转厘米/英寸）"""
    @staticmethod
    def emu_to_cm(emu: float) -> float:
        return round(emu / 360000, 2)  # 1厘米 = 360000 EMU

    @staticmethod
    def emu_to_inch(emu: float) -> float:
        return round(emu / 914400, 2)  # 1英寸 = 914400 EMU

class BaseShape:
    """形状基类（封装位置、大小等通用属性）"""
    def __init__(self, shape):
        self.shape = shape
        self.left = UnitConverter.emu_to_cm(shape.left)
        self.top = UnitConverter.emu_to_cm(shape.top)
        self.width = UnitConverter.emu_to_cm(shape.width)
        self.height = UnitConverter.emu_to_cm(shape.height)

    def get_position(self) -> tuple[float, float]:
        return (self.left, self.top)

    def get_size(self) -> tuple[float, float]:
        return (self.width, self.height)

    def get_box(self) -> tuple[float, float, float, float]:
        return (self.left, self.top, self.width, self.height)

    def get_position_str(self) -> str:
        return f"({self.left},{self.top})"

    def get_size_str(self) -> str:
        return f"({self.width},{self.height})"

class TextBox(BaseShape):
    """文本框解析类"""
    def __init__(self, shape):
        super().__init__(shape)
        self.text_frame = shape.text_frame
        self.text_content = self._extract_text()
        self.paragraph_details = self._extract_paragraph_details()

    def _extract_text(self) -> str:
        return self.text_frame.text.strip() if self.text_frame else ""

    def _extract_paragraph_details(self) -> List[Dict]:
        details = []
        if not self.text_frame:
            return details
        for para in self.text_frame.paragraphs:
            for run in para.runs:
                text = run.text.strip("\n")
                if not text:
                    continue
                details.append({
                    "text": text,
                    "font_name": run.font.name or "默认字体",
                    "font_size": f"{run.font.size.pt:.1f}pt" if run.font.size else "-1pt",
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": self._get_font_color(run.font.color),
                    "alignment": str(para.alignment) if para.alignment else "LEFT (1)"
                })
        return details

    @staticmethod
    def _get_font_color(color) -> str:
        if not color:
            return "无颜色"
        if isinstance(color, RGBColor):
            return f"RGB({color.rgb[0]}, {color.rgb[1]}, {color.rgb[2]})"
        if hasattr(color, "theme_color"):
            return f"主题色: {color.theme_color.name}"
        return "未知颜色"

    def to_dict(self) -> Dict:
        return {
            "type": "文本框",
            "text": self.text_content,
            "position": self.get_position_str(),
            "size": self.get_size_str(),
            "box": self.get_box(),
            "paragraphs": self.paragraph_details
        }

class Table(BaseShape):
    """表格解析类"""
    def __init__(self, shape):
        super().__init__(shape)
        self.table = shape.table
        self.rows = len(self.table.rows)
        self.cols = len(self.table.columns)
        self.cell_details = self._extract_cell_details()
        self.last_row_data = self._extract_last_cell_details()

    @staticmethod
    def _is_row_empty(row, row_idx=0) -> bool:
        """判断表格行是否为空（所有单元格无有效内容）"""
        if row_idx > 0:
            col_idx = 0
            for cell in row.cells:
                stripped_text = cell.text.strip()
                if col_idx > 0:
                    if stripped_text:  # 存在非空白内容
                        return False
                col_idx += 1
            return True
        else:
            for cell in row.cells:
                stripped_text = cell.text.strip()
                if stripped_text:  # 存在非空白内容
                    return False
            return True

    def _extract_cell_details(self) -> List[List[Dict]]:
        details = []
        for row_idx, row in enumerate(self.table.rows):
            row_details = []
            for col_idx, cell in enumerate(row.cells):
                cell_text = "\n".join([para.text.strip() for para in cell.text_frame.paragraphs if para.text.strip()])
                row_details.append({
                    "row": row_idx + 1,
                    "col": col_idx + 1,
                    "text": cell_text,
                    "paragraphs": self._extract_cell_paragraphs(cell)
                })
            details.append(row_details)
        return details

    def _extract_last_cell_details(self) -> list:
        """处理单张表格，返回字典列表（键为标题行，值为数据行）"""
        if len(self.table.rows) < 1:
            return None  # 无行，跳过

        # 提取标题行（首行）
        header_cells = self.table.rows[0].cells
        headers = [cell.text.strip() for cell in header_cells]
        if not any(headers):  # 标题行全为空，无效
            return None

        # 遍历数据行（首行之后的行）
        data = dict.fromkeys(headers)
        for idx in range(1, len(self.table.rows)):
            if self._is_row_empty(self.table.rows[idx], idx):
                continue  # 跳过空行

            # 检查单元格数量是否与标题一致（避免列错位）
            if len(self.table.rows[idx].cells) != len(headers):
                continue  # 列数不匹配，跳过（可根据需求调整）

            # 构建数据字典（标题为键，单元格内容为值）
            row_data = {
                headers[i]: cell.text.strip() for i, cell in enumerate(self.table.rows[idx].cells)
            }
            data = row_data

        return data

    @staticmethod
    def _extract_cell_paragraphs(cell) -> List[Dict]:
        paragraphs = []
        for para in cell.text_frame.paragraphs:
            paragraphs.append({
                "alignment": str(para.alignment) if para.alignment else "Left (1)",
                "runs": [{
                    "text": run.text.strip(),
                    "font_name": run.font.name or "默认字体",
                    "font_size": f"{run.font.size.pt:.1f}pt" if run.font.size else "未设置"
                } for run in para.runs]
            })
        return paragraphs

    def to_dict(self) -> Dict:
        return {
            "type": "Table",
            "box": self.get_box(),
            "position": self.get_position_str(),
            "size": self.get_size_str(),
            "rows": self.rows,
            "cols": self.cols,
            "cells": self.cell_details,
            "last_row": self.last_row_data
        }

class Image(BaseShape):
    """图片解析类"""
    def __init__(self, shape):
        super().__init__(shape)
        self.image = shape
        self.alt_text = self.image.name

    def to_dict(self) -> Dict:
        return {
            "type": "Image",
            "box": self.get_box(),
            "alt_text": self.alt_text,
            "position": self.get_position_str(),
            "size": self.get_size_str()
        }



class CustomShape(BaseShape):
    """AutoShape自定义形状解析类（矩形、圆形等）"""

    def __init__(self, shape):
        super().__init__(shape)
        self.shape_type = self._get_shape_type()
         # 类型标准化：如果是枚举类型，优先用 .name
        shape_type_str = str(self.shape_type)
        if hasattr(self.shape_type, "name"):
            shape_type_str = self.shape_type.name
        if shape_type_str == "RECTANGLE" or shape_type_str == "1":
            self.shape_type = "矩形"
        else:
            self.shape_type = shape_type_str
        self.fill_color = self._get_fill_color()
        self.line_color = self._get_line_color()
        self.line_width = UnitConverter.emu_to_cm(shape.line.width) if shape.line else "无"
        self.text_content = self._get_content()
        self.text_content_font_sizes = []
        self.paragraph_details = self._extract_paragraph_details()

    def _get_content(self):
        if self.shape.has_text_frame:
            return self.shape.text.strip() if len(self.shape.text) > 0 else ""
        else:
            return ""

    def validate_font_size(self, paragraph_details, text_content, target_size=27.0):
        # -------------------- 步驟1：處理文本比對（含首尾空格） --------------------
        # 拚接所有字體塊的text，並去除首尾空格
        full_text = ''.join(item.get('text', '') for item in paragraph_details).strip()
        # 目標text_content也去除首尾空格
        target_text = text_content.strip()

        # 若處理後的文本不一致，直接返回False
        if full_text != target_text:
            return False

        # -------------------- 步驟2：處理字體大小比對（兼容不同格式） --------------------
        # 目標字體大小轉換為浮點數（27.0pt → 27.0；27 pt → 27.0）
        for item in paragraph_details:
            current_text = item.get('text', '')
            if len(current_text) == 0:  # 跳過無實際字元的塊（如分隔符）
                continue

            # 提取當前字體塊的font_size（可能含空格或不同分隔符）
            font_size_str = item.get('font_size', '').strip()  # 先去除前後空格

            # 用正則表達式提取字體數值部分（支援整數/小數，如"27"、"27.0"、"27.5"）
            # 正則說明：\d+ 匹配整數部分，\.?\d* 匹配可選的小數點和小數部分
            match = re.match(r'^(\d+\.?\d*)', font_size_str)
            if not match:
                return False  # 無法識別的字體大小格式，直接失敗

            # 轉換為浮點數（如"27.0"→27.0；"27"→27.0；"27.5"→27.5）
            try:
                current_size = float(match.group(1))
            except ValueError:
                return False  # 數值轉換失敗（理論上不會觸發，因正則已過濾）

            # 比對是否等於目標大小（27.0）
            if not (current_size == target_size):
                return False

        # 所有檢查通過
        return True
    def _get_shape_type(self):
        try:
            # 只对 AutoShape 类型访问 auto_shape_type
            MSO_SHAPE_TYPE.AUTO_SHAPE
            if self.shape.shape_type == 1:  # MSO_SHAPE_TYPE.AUTO_SHAPE
                return getattr(self.shape, "auto_shape_type", "")
            else:
                return str(self.shape.shape_type)
        except Exception:
            print(f"exception : {str(self.shape.shape_type)}")
            return str(self.shape.shape_type)

    def _get_fill_color(self) -> str:
        """解析填充颜色"""
        fill = self.shape.fill
        if fill.type == 1:  # 纯色填充
            if hasattr(fill.fore_color, "rgb"):
                return f"RGB({fill.fore_color.rgb[0]}, {fill.fore_color.rgb[1]}, {fill.fore_color.rgb[2]})"
            if hasattr(fill.fore_color, "theme_color"):
                return f"主题色: {fill.fore_color.theme_color.name}"
        return "无填充色" if fill.type == 0 else "其他填充类型"

    def _get_line_color(self) -> str:
        """解析线条颜色"""
        line = self.shape.line
        if not line:
            return "无线条"
        if hasattr(line.color, "rgb"):
            return f"RGB({line.color.rgb[0]}, {line.color.rgb[1]}, {line.color.rgb[2]})"
        if hasattr(line.color, "theme_color"):
            return f"主题色: {line.color.theme_color.name}"
        return "无线条颜色"

    def _extract_paragraph_details(self) -> list:
        """提取段落级格式和文本片段"""
        details = []
        if not self.shape.has_text_frame:
            return details

        for para in self.shape.text_frame.paragraphs:
            prev_run_info = None
            for run in para.runs:
                text = run.text.strip("\n")
                if not text:
                    continue
                run_info = {
                    "font_name": run.font.name or "默认字体",
                    "font_size": f"{run.font.size.pt:.1f}pt" if run.font.size else "-1pt", # 未設置 默認 -1
                    "bold": run.font.bold,
                    "italic": run.font.italic,
                    "color": self._get_font_color(run.font.color),
                    "alignment": str(para.alignment) if para.alignment else "LEFT (1)"
                }
                if run_info["font_size"] not in self.text_content_font_sizes:
                    self.text_content_font_sizes.append(run_info["font_size"])
                if prev_run_info and all(run_info[k] == prev_run_info[k] for k in run_info):
                    details[-1]["text"] += text
                else:
                    details.append({"text": text, **run_info})
                    prev_run_info = run_info

        # 返回格式化后的列表
        return [
            {"text": item["text"], **{k: v for k, v in item.items() if k != "text"}}
            for item in details
        ]

    @staticmethod
    def _get_font_color(color) -> str:
        """解析字体颜色（支持 RGB 和主题色）"""
        if not color:
            return "无颜色"
        if hasattr(color, "rgb"):
            return f"RGB({color.rgb[0]}, {color.rgb[1]}, {color.rgb[2]})"
        if hasattr(color, "theme_color"):
            return f"主题色: {color.theme_color.name}"
        return "未知颜色"

    def to_string(self) -> str:
        """结构化输出自定义形状信息"""
        return str(self.to_dict())

    def to_dict(self) -> dict:
        return {
            "type": f"{self.shape_type}",
            "box": self.get_box(),
            "text": f"{self.text_content}",
            "position": f"{self.get_position_str()}",
            "size": f"{self.get_size_str()}",
            "fill_color": f"{self.fill_color}",
            "line_color": f"{self.line_color}",
            "line_width": f"{self.line_width:.2f}cm"
        }


# ------------------------------ Slide 类 ------------------------------
class Slide:
    """幻灯片解析类（动态加载标题位置配置）"""

    def __init__(self, slide, slide_master, page_number: int, config_loader: ConfigLoader,
                 version: Optional[str] = None):
        self.slide = slide
        self.page_number = page_number  # 页码从1开始
        if page_number == 7:
            print()
        self.slide_master = slide_master
        self.shapes = self._parse_shapes()
        self.master_shapes = self._parse_master_shapes()

        # 动态加载标题位置配置（从外部传入 config_loader）
        self.title_position = config_loader.get_title_position(version).get("title")
        self.second_title_position = config_loader.get_title_position(version).get("second_title")

        # 初始化标题字段
        self.title: str = ""  # 主标题
        self.second_title: str = ""  # 副标题

        # 自动提取标题（实例化时触发）
        self.extract_titles()

    # ------------------------------ 标题提取方法 ------------------------------
    def extract_titles(self) -> None:
        """从当前幻灯片提取标题和副标题（基于动态配置的位置）"""
        for shape in self.shapes:
            # 筛选文本框或矩形自定义形状
            is_text_shape = (
                    isinstance(shape, TextBox) or
                    (isinstance(shape, CustomShape) and shape.shape_type == "矩形")
            )
            if not is_text_shape or is_empty(shape.text_content):
                continue

            # 获取当前形状的坐标（左, 上, 宽, 高）
            box = (shape.left, shape.top, shape.width, shape.height)

            if self.page_number == 5:
                print()
            # 判断是否为主标题（基于动态加载的 title_position）
            if self._is_title_box(box, self.title_position) and self._validate_size(shape.text_content, 27.0):
                self.title = shape.text_content.strip()

            # 判断是否为副标题（页码>3时生效，基于动态加载的 second_title_position）
            if (
                    self.page_number > 3 and
                    self._is_title_box(box, self.second_title_position) and
                    self._validate_size(shape.text_content, 20.0)
            ):
                self.second_title = shape.text_content.strip()

    @staticmethod
    def _is_title_box(box: Tuple[float, float, float, float],
                      title_box: Tuple[float, float, float, float],
                      iou_threshold: float = 0.3) -> bool:
        """
        判断给定box是否为标题文字栏（基于IOU）
        box: 当前形状坐标 (left, top, width, height)
        title_box: 预设标题框坐标 (left, top, width, height)
        """
        iou = calculate_iou(box, title_box)
        return iou > iou_threshold

    @staticmethod
    def _is_size_target(s: str, target: float) -> Tuple[bool, str]:
        """
        判断字符串是否表示目标尺寸（如27pt）
        返回：(是否匹配, 原始字符串)
        """
        s_clean = s.strip()
        match = re.match(r'^(\d+\.?\d*)', s_clean)  # 提取数字部分（支持整数/小数）
        if not match:
            return False, s

        try:
            num = float(match.group(1))
        except ValueError:
            return False, s

        return num == target, s


    def _validate_size(self, strings: List[str], target_value: float) -> bool:
        """
        验证字符串列表中是否有且仅有一个匹配目标尺寸的字符串
        返回：False（验证通过）/ True（验证不通过）
        """
        size_list = []
        for s in strings:
            is_target, original = Slide._is_size_target(s, target_value)
            if is_target:
                size_list.append(original)

        if not size_list:
            return True  # 无匹配，验证不通过

        unique_formats = set(size_list)  # 去重后的不同格式
        print(
            f"警告：幻灯片 {self.page_number} 发现{len(size_list)}个表示{target_value}pt的字符串，格式为：{unique_formats}")
        return False  # 存在多个/零个匹配，验证不通过

    # ------------------------------ 原有方法 ------------------------------
    def _parse_shapes(self) -> List[BaseShape]:
        """解析当前页的所有形状（原逻辑）"""
        shapes = []
        for shape in self.slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                shapes.append(TextBox(shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                shapes.append(Table(shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                shapes.append(Image(shape))
            elif shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                shapes.append(CustomShape(shape))  # 假设CustomShape已定义
            elif shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                shapes.append(GroupShape(shape))  # 处理群组形状
        return shapes

    def _parse_master_shapes(self) -> List[BaseShape]:
        """解析母版中的文本框（原逻辑）"""
        master_shapes = []
        for layout in self.slide_master.slide_layouts:
            for shape in layout.shapes:
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                    master_shapes.append(TextBox(shape))
        return master_shapes

    def to_dict(self) -> Dict:
        """结构化输出（包含标题）"""
        return {
            "page_number": self.page_number,
            "title": self.title,
            "second_title": self.second_title,
            "shapes": [shape.to_dict() for shape in self.shapes],
            "master_shapes": [shape.to_dict() for shape in self.master_shapes]
        }

class GroupShape(BaseShape):
    """群组形状，包含多个子 shape"""
    def __init__(self, shape):
        super().__init__(shape)
        self.type = "Group"
        self.shapes = []

        # 递归解析所有子 shape
        for sub_shape in shape.shapes:
            if sub_shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX:
                self.shapes.append(TextBox(sub_shape))
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.TABLE:
                self.shapes.append(Table(sub_shape))
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                self.shapes.append(Image(sub_shape))
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.AUTO_SHAPE:
                self.shapes.append(CustomShape(sub_shape))  # 假设CustomShape已定义
            elif sub_shape.shape_type == MSO_SHAPE_TYPE.GROUP:
                self.shapes.append(GroupShape(sub_shape))  # 处理群组形状

    def to_dict(self):
        return {
            "type": self.type,
            "box": self.get_box(),
            "position": self.get_position_str(),
            "size": self.get_size_str(),
            "shapes": [s.to_dict() for s in self.shapes]
        }

if __name__ == "__main__":
    from config.loader import ConfigLoader
    from pptx import Presentation
    # 测试Slide类
    pptx_path = r"D:\pythonf\26xdemo2.pptx"
    prs = Presentation(pptx_path)
    config_loader = ConfigLoader(config_dir="config")
    
    # 测试单个幻灯片
    slide_obj = Slide(prs.slides[0], prs.slide_master, 1, config_loader, "v1")
    slide_dict = slide_obj.to_dict()
    print("Slide转换结果:")
    print(slide_dict)
