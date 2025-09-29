# src/office_ops/ppt_processor/extractors/extractor_a.py
# 实现需求A：提取指定页面的标题等字段并导出为 发包规范
import os
import yaml
import re
import math
from os import path
from typing import List, Dict
from extractors.base_extrator import BaseExtractor
from content_models import Slide, calculate_iou
from utils.text_utils import split_after_colon  # 确保已导入
from utils.logger import LoggerFactory, LOG_LEVELS
from config.loader import ConfigLoader
from pathlib import Path
from openpyxl import load_workbook
from openpyxl.utils.exceptions import InvalidFileException
from openpyxl.worksheet.worksheet import Worksheet  # 新增：导入 Worksheet 类
from typing import List, Dict, Optional, Union
from utils.exceptions import *
# 获取当前文件的绝对路径的根目录
current_file_path = Path(__file__).resolve()
main_dir = current_file_path.parent.parent  # 项目根目录（ppt_processor/）


class ExtractorExcel(BaseExtractor):
    """
    Excel 文件读取器（面向对象风格）
    支持指定标题行、多工作表查询、灵活匹配列

    属性:
        file_path (str): Excel 文件路径
        sheet_name (Optional[str]): 当前工作表名（默认第一个）
        header_row_num (int): 标题行行号（如第4行填4）
        wb (Optional[Workbook]): openpyxl 工作簿对象
        ws (Optional[Worksheet]): 当前工作表对象
        header_row (List[Optional[str]]): 标题行数据（列名列表）
        match_col_idx (Optional[int]): 匹配列的列索引（从0开始）
    """

    def __init__(
            self,
            config: ConfigLoader
    ):
        """
        初始化 Excel 读取器

        参数:
            config: 配置加载器，包含文件路径、工作表名、标题行等配置
        """
        super().__init__()
        # 优化：只加载一次正则规则配置
        self.fields_config = config.config.get("FIELDS_CONFIG", {}).get("发包规范V2_EXCEL", {}).get("header", {})
        self.config = config.get_all_projects_info()
        self.file_path = os.path.join(main_dir, self.config.get('path', 'examples'), self.config.get('filename', '方案总表.xlsx'))
        self.sheet_name = self.config.get('sheet_name', '方案总表')
        self.header_row_num = self.config.get('header_row_num', 4)

        self.logger = LoggerFactory.create_logger("ExtractorExcel")

        # 初始化工作簿和工作表（延迟加载，避免重复打开）
        self.wb: Optional[load_workbook] = None
        self.ws: Optional[Worksheet] = None

        # 初始化标题相关属性
        self.header_row: List[Optional[str]] = []
        self.match_col_idx: Optional[int] = None

    def _load_workbook(self) -> None:
        """私有方法：加载工作簿（仅在首次使用时打开）"""
        if self.wb is None:
            try:
                if not path.exists(self.file_path):
                    raise FileNotFoundError(file_path=self.file_path)  # 自动触发 WARNING 级别日志
                # 以只读模式加载（适合大文件），禁用数据只读（确保能读取所有行）
                self.wb = load_workbook(
                    filename=self.file_path,
                    read_only=True,
                    data_only=True,
                    keep_links=False  # 避免外部链接干扰
                )
                self.logger.debug(f"成功读取 Excel 文件：{self.file_path}")  # 开发模式下控制台+文件输出，生产模式仅文件（若级别允许）
            except OfficeBaseException as e:
                # 记录带错误代码的日志（自动包含堆栈信息）
                self.logger.error(e)
            except InvalidFileException as e:
                self.logger.error(f"文件格式错误：{self.file_path} 不是有效的 .xlsx 文件")
                raise OfficeBaseException(
                    error_code=ErrorCode.FILE_INVALID_ERROR,
                    message=f"文件格式错误：{str(e)}",
                    level=logging.WARNING
                )
                raise ValueError(f"文件格式错误：{self.file_path} 不是有效的 .xlsx 文件")
            except Exception as e:
                # 通用异常处理（非自定义异常）
                self.logger.error(f"未预期的异常：{str(e)}", exc_info=True)
                raise OfficeBaseException(
                    error_code=ErrorCode.UNKNOWN_ERROR,
                    message=f"系统内部错误：{str(e)}",
                    level=logging.CRITICAL
                )

    def _load_worksheet(self) -> None:
        """私有方法：加载工作表（仅在首次使用时打开）"""
        if self.ws is None:
            self._load_workbook()  # 确保工作簿已加载

            # 检查工作表是否存在
            if self.sheet_name and self.sheet_name not in self.wb.sheetnames:
                raise ValueError(
                    f"工作表 '{self.sheet_name}' 不存在，可用工作表：{self.wb.sheetnames}"
                )

            # 选择工作表（优先使用指定的，否则用第一个）
            self.ws = self.wb[self.sheet_name] if self.sheet_name else self.wb.worksheets[0]

    def _read_header(self) -> None:
        """私有方法：读取标题行并用正则规则替换为标准key"""
        if not isinstance(self.ws, Worksheet):
            self._load_worksheet()  # 确保工作表已加载

        max_row = self.ws.max_row
        if self.header_row_num > max_row or self.header_row_num < 1:
            raise ValueError(
                f"标题行号 {self.header_row_num} 无效，Excel 最大行数：{max_row}"
            )

        # 读取标题行（仅读取指定行）
        header_cells = next(
            self.ws.iter_rows(
                min_row=self.header_row_num,
                max_row=self.header_row_num,
                values_only=True
            ),
            []
        )
        
        # 反转为 {标准key: 正则}，方便遍历
        key_regex_map = {k: re.compile(v) for k, v in self.fields_config.items()}

        # 用正则匹配并替换标题
        new_header_row = []
        for cell_value in header_cells:
            replaced = False
            if cell_value is not None:
                for std_key, regex in key_regex_map.items():
                    if regex.search(str(cell_value)):
                        new_header_row.append(std_key)
                        replaced = True
                        break
            if not replaced:
                new_header_row.append(cell_value)

        # new_header_row 排除非空的項后面的 是否存在Nonetype, 如果有則刪除
        while new_header_row and new_header_row[-1] is None:
            new_header_row.pop()
        # new_header_row 檢測中間是否存在Nonetype, 如果有則刪除
        new_header_row = [col if col is not None else f"列{idx + 1}" for idx, col in enumerate(new_header_row)]
        # new_header_row 每個項是否重複 有則拋出警告
        if len(new_header_row) != len(set(new_header_row)):
            self.logger.warning("标题行存在重复项，请检查 Excel 文件")

        self.header_row = new_header_row

    def set_match_column(self, match_column: str) -> None:
        """
        设置匹配列（如 '立項代码'），并自动计算列索引

        参数:
            match_column: 匹配依据的列名（需与标题行一致）
        """
        self._read_header()  # 确保标题行已读取

        if match_column not in self.header_row:
            raise ValueError(
                f"匹配列 '{match_column}' 不存在，标题行列名：{self.header_row}"
            )

        # 计算列索引（从0开始）
        self.match_col_idx = self.header_row.index(match_column)

    def query_by_column(self, target_value: Union[str, int, float]) -> List[Dict]:
        """
        根据指定列匹配目标值，返回匹配的行数据（字典列表）

        参数:
            target_value: 需要匹配的目标值（如立项代码）

        返回:
            匹配的行数据（字典列表，键为列名，值为对应单元格内容）
        """
        # 延迟加载资源和标题
        self._load_worksheet()
        if not self.header_row:
            self._read_header()
        if self.match_col_idx is None:
            raise ValueError("请先通过 set_match_column 方法设置匹配列")

        # 数据行从标题行的下一行开始
        start_row = self.header_row_num + 1
        matched_rows: List[Dict] = []

        # 遍历数据行
        for row in self.ws.iter_rows(min_row=start_row, values_only=True):
            # 跳过空行（整行都是 None）
            if all(cell is None for cell in row):
                continue

            # 检查当前行是否有足够的列（避免索引越界）
            if self.match_col_idx >= len(row):
                continue  # 当前行无对应列数据，跳过

            # 匹配目标值（兼容数值和字符串类型）
            current_value = row[self.match_col_idx]
            if str(current_value) == str(target_value):
                # 构造行字典（缺失列补 None）
                row_dict = {
                    self.header_row[i]: row[i] if i < len(row) else None
                    for i in range(len(self.header_row))
                }
                matched_rows.append(row_dict)

        return matched_rows

    def close(self) -> None:
        """关闭工作簿（无需检查 is_closed）"""
        if self.wb:
            self.wb.close()  # openpyxl 的 Workbook.close() 无需检查状态
            self.wb = None  # 释放引用
            self.ws = None  # 释放工作表引用

    def __enter__(self):
        """上下文管理器：进入时自动加载资源"""
        self._load_worksheet()
        self._read_header()
        return self

    def __exit__(self, exc_type, exc_val, exc_tb):
        """上下文管理器：退出时自动关闭工作簿"""
        self.close()

    def extract(self) -> List[Dict]:
        """
        结构化返回所有数据行，每行一个dict，header为key，数据为value
        """
        self._load_worksheet()
        self._read_header()
        start_row = self.header_row_num + 1
        result = []
        for row in self.ws.iter_rows(min_row=start_row, values_only=True):
            # 跳过空行
            if all(cell is None for cell in row):
                continue
            # 构造行字典（缺失列补None）
            row_dict = {
                self.header_row[i]: row[i] if i < len(row) else None
                for i in range(len(self.header_row))
            }
            result.append(row_dict)
        return result


class ExtractorA(BaseExtractor):
    """需求A提取器（支持多类型 shape 提取）"""
    def __init__(self, slides: List[Slide], config_loader: ConfigLoader = None):
        self.logger = LoggerFactory.create_logger("Extractor_发包规范")
        self.logger.info("初始化提取器")
        super().__init__(slides)
        # 优先使用传入的config_loader，否则尝试从fields_loader获取
        if config_loader:
            self.config = config_loader.config.get("FIELDS_CONFIG", {}).get("发包规范V1_PPT", {})
        else:
            self.config = {}

        if not self.config:
            self.logger.warning("未找到发包规范V1_PPT的配置")


    def _calc_utilization_rate(self, failure_rate: str) -> str:
            """
            根据故障率字符串自动计算利用率（百分比取反），如 '≤1.23%' -> '≥98.77%'
            """
            import re
            # 去掉所有空格
            failure_rate = failure_rate.replace(" ", "")
            match = re.search(r"([≤<]?)(\d+(\.\d+)?)%", failure_rate)
            if match:
                num = float(match.group(2))
                utilization = 100 - num
                # 保留两位小数
                return f"≥ {utilization:.2f}%"
            return ""

    def extract(self) -> Dict:
        self.logger.info("excel 开始提取内容")
        try:
            flat_result = {}

            # 处理 master 字段（用 master_shapes 匹配）
            master_cfg = self.config.get("master", {})
            if master_cfg:
                pos_dict = master_cfg.get("iou", {})
                iou_threshold = master_cfg.get("iou_threshold", 0.1)
                for key, position in pos_dict.items():
                    if len(position) < 4:
                        self.logger.error(f"position of ProjectCode is empty, please check the fields_config.yaml")
                        raise ValueError("position of ProjectCode is empty, please check the fields_config.yaml")
                    # 遍历所有 slide，查找 master_shapes
                    for slide in self.slides:
                        for shape in slide.get("master_shapes", []):
                            iou = calculate_iou(shape["box"], position)
                            self.logger.debug(f"匹配 master 字段 {key}，计算 IOU: {iou:.2f}")
                            if iou > iou_threshold:
                                # 只取第一个匹配
                                flat_result[key] = shape.get("text", "")
                                break

            # 按 page 定位
            for page_num, page_cfg in self.config.get("page", {}).items():
                slide = self._get_slide_by_page(int(page_num))
                if not slide:
                    continue
                for field, iou_cfg in page_cfg.get("iou", {}).items():
                    box = iou_cfg.get("box")
                    need_split = iou_cfg.get("need_split")
                    storage_vars = iou_cfg.get("storage_var")
                    # IOU 匹配 shape
                    for shape in slide["shapes"]:
                        iou = calculate_iou(shape["box"], box)
                        if iou > 0.3:
                            text = shape.get("text", "")
                            if need_split and isinstance(storage_vars, list):
                                parts = text.split(need_split)
                                for idx, var in enumerate(storage_vars):
                                    if var and idx < len(parts):
                                        flat_result[var] = parts[idx]
                            else:
                                flat_result[field] = text
                            break

            # 按 title 定位
            for title_cfg in self.config.get("title", []):
                slide = self._get_slide_by_title(title_cfg["first"], title_cfg.get("second", ""))
                if not slide:
                    continue
                # re 匹配文本字段
                for field, field_cfg in title_cfg.get("re", {}).items():
                    storage_var = field
                    flat_result[storage_var] = self._extract_text_by_re(slide, field, field_cfg)
                # table 匹配表格字段
                for field, field_cfg in title_cfg.get("table", {}).items():
                    storage_var = field
                    flat_result[storage_var] = self._extract_table_last_row(slide, field_cfg["match_key_string"])

            # 自动补充 dev_utilization_rate 字段
            failure_rate = flat_result.get("dev_failure_rate", "")
            if failure_rate and isinstance(failure_rate, str):
                flat_result["dev_utilization_rate"] = self._calc_utilization_rate(failure_rate)

            self.logger.info("内容提取完成")
            return flat_result
        except Exception as e:
            self.logger.error(f"提取内容时出错: {str(e)}", exc_info=True)

    def _get_slide_by_page(self, page_number: int):
        for slide in self.slides:
            if slide["page_number"] == page_number:
                return slide
        return None

    def _get_slide_by_title(self, first_title: str, second_title: str):
        for slide in self.slides:
            if slide["title"] == first_title and (not second_title or slide["second_title"] == second_title):
                return slide
        return None

    def _extract_shape_by_position(self, slide, position):
        """
        通过 IOU 判断形状位置是否匹配，提取内容
        position: (left, top, width, height)
        """
        iou_threshold = 0.3  # 可根据实际情况调整
        for shape in slide["shapes"]:
            iou = calculate_iou(shape["box"], position)
            if iou > iou_threshold:
                if shape["type"] == "TextBox" or shape["type"] == "文本框" or shape["type"] == "矩形":
                    return shape.get("text", "")
        return None

    def _extract_text_by_re(self, slide, field_name, field_cfg):
        """
        根据正则表达式提取文本内容或相关形状
        
        Args:
            slide: 幻灯片数据
            field_name: 字段名称
            field_cfg: 字段配置
            
        Returns:
            根据match_rule返回不同类型的数据:
            - match_rule > 0: 返回带页码的shape对象
            - match_rule = 0: 返回匹配的文本
            - match_rule = -1: 返回冒号后的提取值
        """
        match_rule = field_cfg.get("match_rule", 0)
        re_rule = field_cfg.get("re_rule", "")
        shapes = slide["shapes"]
        direction_map = {1: "down", 2: "left", 3: "up", 4: "right"}

        # 遍历所有shape
        for shape in shapes:
            # 处理群组
            if shape["type"] == "Group":
                result = self._process_group_shape(
                    shape, re_rule, match_rule, 
                    slide["page_number"], direction_map
                )
                if result:
                    return result
                    
            # 处理独立文本框
            elif shape["type"] in ["文本框", "矩形"]:
                result = self._process_single_shape(
                    shape, re_rule, match_rule,
                    shapes, slide["page_number"], direction_map
                )
                if result:
                    return result

        return None

    def _process_group_shape(self, group, re_rule, match_rule, page_number, direction_map):
        """处理群组内的shape"""
        for sub_shape in group["shapes"]:
            if sub_shape["type"] in ["文本框", "矩形"]:
                text = sub_shape.get("text", "")
                if re_rule and re.search(re_rule, text):
                    if match_rule > 0:
                        if match_rule == 5:
                            # 返回当前sub_shape本身和页码
                            sub_shape["page_number"] = page_number
                            return sub_shape
                        else:
                            # 原有的群组内搜索逻辑
                            for target in group["shapes"]:
                                if target["type"] in ["Image", "CustomShape", "Group", "矩形"]:
                                    target["page_number"] = page_number
                                    return target
                    elif match_rule == 0:
                        # 清理文本前后的空白字符
                        return text.strip() if text else ""
                    elif match_rule == -1:
                        _text = self._extract_value_after_colon(text, re_rule)
                        return _text.strip() if _text else ""
        return None

    def _process_single_shape(self, shape, re_rule, match_rule, all_shapes, page_number, direction_map):
        """处理单独的文本框"""
        text = shape.get("text", "")
        if re_rule and re.search(re_rule, text):
            if match_rule > 0:
                if match_rule == 5:
                    # 返回当前shape本身和页码
                    shape["page_number"] = page_number
                    return shape
                else:
                    # 原有的近邻搜索逻辑
                    direction = direction_map.get(match_rule, "down")
                    target = self._find_nearest_shape(shape, all_shapes, direction)
                    if target:
                        target["page_number"] = page_number
                        return target
            elif match_rule == 0:
                # 清理文本前后的空白字符
                return text.strip() if text else ""
            elif match_rule == -1:
                _text = self._extract_value_after_colon(text, re_rule)
                return _text.strip() if _text else ""
        return None

    def _extract_value_after_colon(self, text: str, re_rule: str) -> str:
        """提取冒号后的值"""
        regex = re.compile(re_rule)
        match = regex.search(text)
        if match and len(match.regs) > 0:
            temp = text[match.regs[0][0]:match.regs[0][1]]
            return split_after_colon(temp) or ""
        return ""

    def _find_nearest_shape(self, ref_shape, shapes, direction="down"):
        """
        查找与 ref_shape 最近的目标 shape（如图片或 custom shape），支持下方/左侧/上方/右侧
        """
        ref_box = ref_shape["box"]
        ref_center = (ref_box[0] + ref_box[2] / 2, ref_box[1] + ref_box[3] / 2)
        min_dist = float("inf")
        nearest = None
        for shape in shapes:
            if shape is ref_shape:
                continue
            # 目标类型可根据需求调整
            if shape["type"] in ["Image", "图片", "矩形", "Group"]:
                box = shape["box"]
                center = (box[0] + box[2] / 2, box[1] + box[3] / 2)
                # 按方向筛选
                if direction == "down" and center[1] > ref_center[1]:
                    dist = math.hypot(center[0] - ref_center[0], center[1] - ref_center[1])
                elif direction == "up" and center[1] < ref_center[1]:
                    dist = math.hypot(center[0] - ref_center[0], ref_center[1] - center[1])
                elif direction == "left" and center[0] < ref_center[0]:
                    dist = math.hypot(ref_center[0] - center[0], center[1] - ref_center[1])
                elif direction == "right" and center[0] > ref_center[0]:
                    dist = math.hypot(center[0] - ref_center[0], center[1] - ref_center[1])
                else:
                    continue
                if dist < min_dist:
                    min_dist = dist
                    nearest = shape
        return nearest

    def _extract_table_last_row(self, slide, field_name):
        for shape in slide["shapes"]:
            if shape["type"] == "Table":
                return shape["last_row"].get(field_name, "")
        return None
    

def example_usage_excel():
    """示例用法"""
    # 配置参数（根据实际情况修改）
    match_column = "立項代码"  # 匹配列名
    target_value = "adazffa5"  # 要匹配的目标值
    config_loader = ConfigLoader(config_dir=r"..\config")

    try:
        # 方式1：手动创建对象并调用方法
        reader = ExtractorExcel(config=config_loader)
        result = reader.extract()
        # 若是已經有了result reader 可以退出了，就不會占用excel資源
        reader.close()
        print(f"总共读取 {len(result)} 行数据")
        # # 方式2：使用上下文管理器（自动管理资源）
        # with ExtractorExcel(
        #         config=config_loader,
        # ) as reader:
        #     reader.set_match_column(match_column)
        #     result = reader.query_by_column(target_value)
        #     print(f"上下文管理器查询结果：{len(result)} 条")

        # 打印结果（示例）
        if result:
            print("\n匹配结果：")
            print(" | ".join(result[0].keys()))  # 打印表头
            for row in result:
                print(" | ".join(str(v) if v is not None else "空" for v in row.values()))

    except ValueError as e:
        print(f"错误：{str(e)}")
    except Exception as e:
        print(f"未知错误：{str(e)}")

def example_usage_extractor_a():
    """示例用法"""
    from config.loader import ConfigLoader
    from content_models import Slide
    from pptx import Presentation
    
    # 测试提取器
    pptx_path = r"..\examples/templates/26xdemo2.pptx"
    config_loader = ConfigLoader(config_dir=r"..\config")


    
    # 读取PPT并结构化
    prs = Presentation(pptx_path)
    slides = []
    for page_number, slide in enumerate(prs.slides, start=1):
        slide_obj = Slide(slide, prs.slide_master, page_number, config_loader, "v1")
        slides.append(slide_obj)
    
    slides_dicts = [slide.to_dict() for slide in slides]
    
    # 测试提取
    extractor = ExtractorA(slides_dicts, config_loader)
    result = extractor.extract()
    print("字段提取完成")
    print(f"提取结果: {result}")


if __name__ == "__main__":
    # example_usage_excel()
    example_usage_extractor_a()
    pass